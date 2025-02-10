import streamlit as st
from langchain_community.document_loaders import PDFPlumberLoader
from langchain_community.document_loaders import Docx2txtLoader
import pandas as pd
from langchain.docstore.document import Document
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_ollama import ChatOllama
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
import tempfile
import os
import shutil
from chromadb.config import Settings
import pptx

# 设置页面标题
st.set_page_config(page_title="文件问答助手", layout="wide")
st.title("文件问答助手")

# 初始化会话状态
if "messages" not in st.session_state:
    st.session_state.messages = []
if "vectorstore" not in st.session_state:
    st.session_state.vectorstore = None

# 创建持久化目录
PERSIST_DIR = "./chroma_db"
if not os.path.exists(PERSIST_DIR):
    os.makedirs(PERSIST_DIR)

# 添加 Chroma 设置
CHROMA_SETTINGS = Settings(
    allow_reset=True,
    is_persistent=True,
    persist_directory=PERSIST_DIR
)

# 文件上传部分
uploaded_file = st.file_uploader("上传文件", type=['pdf', 'xlsx', 'xls', 'docx', 'doc', 'pptx', 'ppt'])

if uploaded_file is not None:
    # 创建临时文件保存上传的文件
    file_extension = uploaded_file.name.split('.')[-1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    try:
        # 根据文件类型选择合适的加载器
        if file_extension in ['pdf']:
            loader = PDFPlumberLoader(tmp_file_path)
            docs = loader.load()
        elif file_extension in ['xlsx', 'xls']:
            # 使用 pandas 读取 Excel
            df = pd.read_excel(tmp_file_path)
            docs = []
            for index, row in df.iterrows():
                content = " ".join(str(cell) for cell in row)
                metadata = {"source": f"row_{index}", "file": uploaded_file.name}
                docs.append(Document(page_content=content, metadata=metadata))
        elif file_extension in ['docx', 'doc']:
            loader = Docx2txtLoader(tmp_file_path)
            docs = loader.load()
        elif file_extension in ['pptx', 'ppt']:
            # 使用 python-pptx 处理 PPT
            prs = pptx.Presentation(tmp_file_path)
            docs = []
            for slide_number, slide in enumerate(prs.slides, 1):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                
                if text_content:
                    content = "\n".join(text_content)
                    metadata = {
                        "source": f"slide_{slide_number}",
                        "file": uploaded_file.name
                    }
                    docs.append(Document(page_content=content, metadata=metadata))
        
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=0)
        all_splits = text_splitter.split_documents(docs)
        
        # 创建向量存储
        local_embeddings = OllamaEmbeddings(model="nomic-embed-text")
        st.session_state.vectorstore = Chroma.from_documents(
            documents=all_splits, 
            embedding=local_embeddings,
            collection_name="document_collection",
            persist_directory=PERSIST_DIR,
            client_settings=CHROMA_SETTINGS
        )
        
        # 清理临时文件
        os.unlink(tmp_file_path)
        st.success(f"{uploaded_file.name} 已成功上传并处理！")
        
    except Exception as e:
        st.error(f"处理文件时发生错误: {str(e)}")
        os.unlink(tmp_file_path)

# 设置模型和提示模板
model = ChatOllama(model="deepseek-r1:7b")

SYSTEM_TEMPLATE = """
Role: 文件信息处理专家

Profile
Language: 中文
Description: 我是专业的文件内容解析助手，擅长从用户提供的文件资料中精准提取信息，结合上下文进行深度推理分析，提供准确可靠的问答服务。

Skill-1 信息检索
1. 运用RAG技术从上传文件中提取关键信息
2. 支持PDF/DOCX/XLSX/PPTX等多种格式文件解析
3. 自动识别文档结构和语义关系

Rules
0. 用户所提到的一切文件，包括PDF、Excel、Word、PPT等，都是指上下文内容
1. 严格基于提供的上下文回答问题
2. 未经明确允许绝不假设或推测信息
3. 遇到超出知识范围的问题明确告知
4. 保持专业客观，不添加主观意见
5. 不讨论与文件内容无关的话题

Workflow
1. 接收并解析用户上传的文件内容
2. 分析问题需求与上下文关联性
3. 从检索到的信息中提取关键证据
4. 组织逻辑清晰的完整详细的回答
"""

RAG_TEMPLATE = """
{system_message}

请使用以下检索到的上下文来回答问题。如果你不知道答案，请直接说不知道。

<context>
{context}
</context>

问题: {question}

回答要求：
1. 如上下文无相关信息请直接说明
2. 关键数据需标注来源段落
3. 保持回答详细专业
"""

rag_prompt = ChatPromptTemplate.from_messages([
    ("system", SYSTEM_TEMPLATE),
    ("user", RAG_TEMPLATE)
])

# 格式化文档的函数
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# 显示聊天历史
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# 聊天输入
if question := st.chat_input("请输入您的问题"):
    # 添加用户问题到聊天历史
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)

    # 检查是否已上传文件
    if st.session_state.vectorstore is None:
        with st.chat_message("assistant"):
            st.markdown("请先上传文件！")
        st.session_state.messages.append({"role": "assistant", "content": "请先上传文件！"})
    else:
        # 创建检索链
        retriever = st.session_state.vectorstore.as_retriever()
        qa_chain = (
            {
                "context": retriever | format_docs, 
                "question": lambda x: x,
                "system_message": lambda _: SYSTEM_TEMPLATE
            }
            | rag_prompt
            | model
            | StrOutputParser()
        )

        # 显示助手回答
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            full_response = ""
            
            # 使用流式输出
            for chunk in qa_chain.stream(question):
                full_response += chunk
                message_placeholder.markdown(full_response + "▌")
            
            # 完成后显示完整消息，去掉光标
            message_placeholder.markdown(full_response)
            # 将完整回答添加到聊天历史
            st.session_state.messages.append({"role": "assistant", "content": full_response})

# 修改清除聊天历史的按钮处理逻辑
if st.button("清除聊天历史"):
    # 清除聊天消息
    st.session_state.messages = []
    
    # 清除向量存储
    if st.session_state.vectorstore is not None:
        try:
            # 确保关闭当前的 client 连接
            if hasattr(st.session_state.vectorstore._client, '_producer'):
                st.session_state.vectorstore._client._producer.close()
            
            # 删除持久化目录
            if os.path.exists(PERSIST_DIR):
                try:
                    shutil.rmtree(PERSIST_DIR)
                except PermissionError:
                    # 如果遇到权限错误，先修改文件权限
                    for root, dirs, files in os.walk(PERSIST_DIR):
                        for dir in dirs:
                            os.chmod(os.path.join(root, dir), 0o777)
                        for file in files:
                            os.chmod(os.path.join(root, file), 0o777)
                    shutil.rmtree(PERSIST_DIR)
                
                # 重新创建空目录
                os.makedirs(PERSIST_DIR)
            
            # 清除向量存储的会话状态
            st.session_state.vectorstore = None
            
        except Exception as e:
            st.error(f"清理数据时发生错误: {str(e)}")
            # 尝试强制重置状态
            st.session_state.vectorstore = None
    
    st.rerun()